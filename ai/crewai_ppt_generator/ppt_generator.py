import os

import matplotlib
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
matplotlib.use("Agg")  # Use non-interactive backend
# --- 한글 폰트 설정 시작 ---
import os
FONT_PATHS = [
    "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",  # Ubuntu/Debian
    "/usr/share/fonts/truetype/nanum/NanumGothic-Regular.ttf",  # 일부 시스템
    "/usr/share/fonts/truetype/malgun/MalgunGothic.ttf",  # Windows WSL
    "/Library/Fonts/AppleGothic.ttf",  # macOS
]
font_found = False
for font_path in FONT_PATHS:
    if os.path.exists(font_path):
        try:
            fm.fontManager.addfont(font_path)
            font_prop = fm.FontProperties(fname=font_path)
            matplotlib.rc("font", family=font_prop.get_name())
            font_found = True
            break
        except Exception:
            continue
if not font_found:
    print("경고: 한글 폰트를 찾을 수 없습니다. 시스템에 NanumGothic 또는 대체 폰트를 설치하세요.")
matplotlib.rcParams["axes.unicode_minus"] = False
# --- 한글 폰트 설정 끝 ---
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


from typing import Optional, Any, Dict, List, Union
import pandas as pd
from pptx import Presentation

class PPTGenerator:
    """
    Class for generating PowerPoint presentations from agent-generated content
    """

    def __init__(self, report_topic: str, output_dir: str = "output") -> None:
        """
        Initialize the PPT generator

        Args:
            report_topic (str): The main topic of the report
            output_dir (str): Directory to save the output files
        """
        self.report_topic = report_topic
        self.output_dir = output_dir
        self.prs = Presentation()
        self.chart_count = 0

        # Create output directory if it doesn't exist
        os.makedirs(output_dir, exist_ok=True)

        # Set up slide layouts for easy reference
        self.title_slide_layout = self.prs.slide_layouts[0]  # Title Slide
        self.title_content_layout = self.prs.slide_layouts[1]  # Title and Content
        self.section_layout = self.prs.slide_layouts[2]  # Section Header
        self.two_content_layout = self.prs.slide_layouts[3]  # Two Content
        self.title_only_layout = self.prs.slide_layouts[5]  # Title Only
        self.blank_layout = self.prs.slide_layouts[6]  # Blank

    def add_title_slide(self, subtitle: str = "Generated by CrewAI Multi-Agent System") -> None:
        """
        Add a title slide to the presentation

        Args:
            subtitle (str): Subtitle for the title slide
        """
        slide = self.prs.slides.add_slide(self.title_slide_layout)
        title = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        title.text = f"{self.report_topic}"
        subtitle_shape.text = subtitle

        # Format the title
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    def add_section_slide(self, section_title: str) -> None:
        """
        Add a section slide to the presentation

        Args:
            section_title (str): Title for the section
        """
        slide = self.prs.slides.add_slide(self.section_layout)
        title = slide.shapes.title
        title.text = section_title

        # Format the title
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    def add_content_slide(self, title: str, content: str) -> None:
        """
        Add a content slide to the presentation

        Args:
            title (str): Title for the slide
            content (str): Text content for the slide
        """
        slide = self.prs.slides.add_slide(self.title_content_layout)
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]

        title_shape.text = title
        content_shape.text = content

        # Format the title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

        # Format the content
        for paragraph in content_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(18)

    def add_two_column_slide(self, title: str, left_content: str, right_content: str) -> None:
        """
        Add a slide with two columns of content

        Args:
            title (str): Title for the slide
            left_content (str): Content for the left column
            right_content (str): Content for the right column
        """
        slide = self.prs.slides.add_slide(self.two_content_layout)
        title_shape = slide.shapes.title
        left_shape = slide.placeholders[1]
        right_shape = slide.placeholders[2]

        title_shape.text = title
        left_shape.text = left_content
        right_shape.text = right_content

        # Format the title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

        # Format the content
        for shape in [left_shape, right_shape]:
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.size = Pt(18)

    def create_chart(self, chart_type: str, data: Union[dict, pd.DataFrame], title: str, x_label: Optional[str] = None, y_label: Optional[str] = None) -> str:
        """
        Create a chart image file

        Args:
            chart_type (str): Type of chart ('bar', 'line', 'pie', etc.)
            data (dict or pd.DataFrame): Data for the chart
            title (str): Title for the chart
            x_label (str, optional): Label for the x-axis
            y_label (str, optional): Label for the y-axis

        Returns:
            str: Path to the saved chart image
        """
        self.chart_count += 1
        plt.figure(figsize=(10, 6))

        if chart_type == "bar":
            if isinstance(data, dict):
                plt.bar(list(data.keys()), list(data.values()))
            else:  # Assume DataFrame
                data.plot(kind="bar", ax=plt.gca())

        elif chart_type == "line":
            if isinstance(data, dict):
                plt.plot(list(data.keys()), list(data.values()))
            else:  # Assume DataFrame
                data.plot(kind="line", ax=plt.gca())

        elif chart_type == "pie":
            if isinstance(data, dict):
                plt.pie(
                    list(data.values()), labels=list(data.keys()), autopct="%1.1f%%"
                )
            else:  # Assume DataFrame
                data.plot(kind="pie", ax=plt.gca())

        elif chart_type == "scatter":
            if isinstance(data, dict):
                # Assume dict with 'x' and 'y' keys
                plt.scatter(data["x"], data["y"])
            else:  # Assume DataFrame
                data.plot(
                    kind="scatter", x=data.columns[0], y=data.columns[1], ax=plt.gca()
                )

        # Set title and labels
        plt.title(title, fontsize=14, fontweight="bold")
        if x_label:
            plt.xlabel(x_label)
        if y_label:
            plt.ylabel(y_label)

        # Add grid and style
        plt.grid(True, linestyle="--", alpha=0.7)
        plt.tight_layout()

        # Save the chart
        chart_path = os.path.join(self.output_dir, f"chart_{self.chart_count}.png")
        plt.savefig(chart_path, dpi=300, bbox_inches="tight")
        plt.close()

        return chart_path

    def add_chart_slide(self, title: str, chart_path: str, description: Optional[str] = None) -> None:
        """
        Add a slide with a chart

        Args:
            title (str): Title for the slide
            chart_path (str): Path to the chart image
            description (str, optional): Description text for the chart
        """
        slide = self.prs.slides.add_slide(self.title_only_layout)
        title_shape = slide.shapes.title
        title_shape.text = title

        # Format the title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

        # Add the chart
        left = Inches(1.5)
        top = Inches(1.8)
        width = Inches(7)
        height = Inches(4)
        slide.shapes.add_picture(chart_path, left, top, width, height)

        # Add description if provided
        if description:
            left = Inches(1)
            top = Inches(6)
            width = Inches(8)
            height = Inches(1)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = description

            # Format the description
            for paragraph in tf.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.italic = True

    def parse_content_sections(self, content: str) -> dict[str, str]:
        """
        Parse content into sections based on headings

        Args:
            content (str): The full content text

        Returns:
            dict: Dictionary of section titles and their content
        """
        # Print the content for debugging
        print(f"Content to parse: {content[:200]}...")

        # Define patterns for section headings (both English and Korean)
        section_patterns = [
            # English patterns
            r"Executive Summary[:\s]*(.*?)(?=\n\s*\n\s*[A-Z가-힣][\w\s가-힣]+[:\n]|$)",
            r"Introduction[:\s]*(.*?)(?=\n\s*\n\s*[A-Z가-힣][\w\s가-힣]+[:\n]|$)",
            r"Key Findings[:\s]*(.*?)(?=\n\s*\n\s*[A-Z가-힣][\w\s가-힣]+[:\n]|$)",
            r"Strategic Recommendations[:\s]*(.*?)(?=\n\s*\n\s*[A-Z가-힣][\w\s가-힣]+[:\n]|$)",
            r"Conclusion[:\s]*(.*?)(?=\n\s*\n\s*[A-Z가-힣][\w\s가-힣]+[:\n]|$)",
            # Korean patterns
            r"임원 요약[:\s]*(.*?)(?=\n\s*\n\s*[A-Z가-힣][\w\s가-힣]+[:\n]|$)",
            r"소개[:\s]*(.*?)(?=\n\s*\n\s*[A-Z가-힣][\w\s가-힣]+[:\n]|$)",
            r"주요 발견[:\s]*(.*?)(?=\n\s*\n\s*[A-Z가-힣][\w\s가-힣]+[:\n]|$)",
            r"전략적 권장 사항[:\s]*(.*?)(?=\n\s*\n\s*[A-Z가-힣][\w\s가-힣]+[:\n]|$)",
            r"결론[:\s]*(.*?)(?=\n\s*\n\s*[A-Z가-힣][\w\s가-힣]+[:\n]|$)",
        ]

        # Both English and Korean section titles
        section_titles = [
            "Executive Summary",
            "Introduction",
            "Key Findings",
            "Strategic Recommendations",
            "Conclusion",
            # Map Korean titles to English section names for consistency
            "Executive Summary",
            "Introduction",
            "Key Findings",
            "Strategic Recommendations",
            "Conclusion",
        ]

        sections = {}

        # Extract content for each section
        for i, pattern in enumerate(section_patterns):
            match = re.search(pattern, content, re.DOTALL)
            if match:
                section_key = section_titles[
                    i % 5
                ]  # Use modulo to map to the 5 main sections
                sections[section_key] = match.group(1).strip()
                print(
                    f"Found section {section_key} with content length: {len(sections[section_key])}"
                )

        # If no sections were found, try to extract content using headings
        if not any(sections.values()):
            print(
                "No sections found with regex patterns, trying heading-based extraction"
            )
            # Split by common heading markers
            lines = content.split("\n")
            current_section = None
            section_content = []

            for line in lines:
                # Check if this line looks like a heading
                if re.match(
                    r"^#+\s+|^\d+\.\s+|^[A-Z가-힣][A-Z가-힣\s]+:$", line.strip()
                ):
                    # Save previous section if we were collecting one
                    if current_section and section_content:
                        sections[current_section] = "\n".join(section_content)
                        section_content = []

                    # Determine which section this heading belongs to
                    lower_line = line.lower()
                    if any(
                        term in lower_line
                        for term in ["executive", "summary", "요약", "개요"]
                    ):
                        current_section = "Executive Summary"
                    elif any(
                        term in lower_line for term in ["introduction", "소개", "서론"]
                    ):
                        current_section = "Introduction"
                    elif any(
                        term in lower_line
                        for term in ["finding", "analysis", "발견", "분석"]
                    ):
                        current_section = "Key Findings"
                    elif any(
                        term in lower_line
                        for term in ["recommendation", "strategy", "권장", "전략"]
                    ):
                        current_section = "Strategic Recommendations"
                    elif any(
                        term in lower_line
                        for term in ["conclusion", "summary", "결론", "요약"]
                    ):
                        current_section = "Conclusion"
                    else:
                        current_section = None
                elif current_section:
                    section_content.append(line)

            # Save the last section
            if current_section and section_content:
                sections[current_section] = "\n".join(section_content)

        # If still no sections found, create a default structure from the content
        if not any(sections.values()):
            print("Creating default sections from content")
            content_parts = content.split("\n\n")
            total_parts = len(content_parts)

            if total_parts >= 5:
                sections["Executive Summary"] = content_parts[0]
                sections["Introduction"] = content_parts[1]
                sections["Key Findings"] = "\n\n".join(
                    content_parts[2 : total_parts - 2]
                )
                sections["Strategic Recommendations"] = content_parts[total_parts - 2]
                sections["Conclusion"] = content_parts[total_parts - 1]
            elif total_parts >= 3:
                sections["Executive Summary"] = content_parts[0]
                sections["Key Findings"] = "\n\n".join(
                    content_parts[1 : total_parts - 1]
                )
                sections["Conclusion"] = content_parts[total_parts - 1]
            else:
                # Just put all content in Key Findings if we can't split it
                sections["Key Findings"] = content

        return sections

    def generate_from_agent_results(self, data_analysis: str, content: str, visual_design: str) -> str:
        """
        Generate a complete presentation from agent results

        Args:
            data_analysis (str): Results from the data analysis agent
            content (str): Results from the content writing agent
            visual_design (str): Results from the visual design agent

        Returns:
            str: Path to the generated PowerPoint file
        """
        # Parse content into sections
        sections = self.parse_content_sections(content)

        # Add title slide
        self.add_title_slide()

        # Add executive summary
        if sections.get("Executive Summary"):
            self.add_section_slide("Executive Summary")
            self.add_content_slide("Key Points", sections["Executive Summary"])

        # Add introduction
        if sections.get("Introduction"):
            self.add_section_slide("Introduction")
            self.add_content_slide("Background", sections["Introduction"])

        # Extract chart suggestions from data analysis
        chart_suggestions = self.extract_chart_suggestions(data_analysis)

        # Add key findings with charts
        if sections.get("Key Findings"):
            self.add_section_slide("Key Findings")

            # Split key findings into paragraphs
            findings = sections["Key Findings"].split("\n\n")

            # Add each finding as a separate slide
            for i, finding in enumerate(findings):
                if finding.strip():
                    # Add chart if available
                    if i < len(chart_suggestions):
                        chart_info = chart_suggestions[i]
                        chart_path = self.create_chart(
                            chart_info["type"],
                            chart_info["data"],
                            chart_info["title"],
                            chart_info.get("x_label"),
                            chart_info.get("y_label"),
                        )
                        self.add_chart_slide(f"Finding {i+1}", chart_path, finding)
                    else:
                        self.add_content_slide(f"Finding {i+1}", finding)

        # Add strategic recommendations
        if sections.get("Strategic Recommendations"):
            self.add_section_slide("Strategic Recommendations")
            self.add_content_slide(
                "Recommendations", sections["Strategic Recommendations"]
            )

        # Add conclusion
        if sections.get("Conclusion"):
            self.add_section_slide("Conclusion")
            self.add_content_slide("Summary", sections["Conclusion"])

        # Save the presentation
        output_file = os.path.join(
            self.output_dir, f"{self.report_topic.replace(' ', '_')}_Report.pptx"
        )
        self.prs.save(output_file)

        return output_file

    def extract_chart_suggestions(self, data_analysis: str) -> list[dict[str, Any]]:
        """
        Extract chart suggestions from data analysis text

        Args:
            data_analysis (str): Text from data analysis agent

        Returns:
            list: List of chart information dictionaries
        """
        print(
            f"Extracting chart suggestions from data analysis: {data_analysis[:200]}..."
        )

        chart_suggestions = []

        # Try to extract actual data from the analysis text
        # Look for patterns that might indicate data points

        # Pattern 1: Look for percentage values (e.g., 45%, 20%, 18%)
        percentage_matches = re.findall(
            r"(\d+)\s*%\s*(?:의|of)?\s*([^,.\n]+)", data_analysis
        )
        if percentage_matches and len(percentage_matches) >= 3:
            # Create a pie chart from percentages
            percentage_data = {}
            for match in percentage_matches[:5]:  # Take up to 5 data points
                value = int(match[0])
                label = match[1].strip()
                percentage_data[label] = value

            chart_suggestions.append(
                {
                    "type": "pie",
                    "title": "분포 분석",  # Distribution Analysis
                    "data": percentage_data,
                }
            )

        # Pattern 2: Look for quarterly data
        quarters = ["Q1", "Q2", "Q3", "Q4", "1분기", "2분기", "3분기", "4분기"]
        quarter_data = {}

        for quarter in quarters:
            # Look for patterns like "Q1: 120,000" or "1분기: 120,000"
            matches = re.findall(f"{quarter}[^\d]+(\d[\d,.]+)", data_analysis)
            if matches:
                # Clean the number and convert to float
                value_str = matches[0].replace(",", "")
                try:
                    quarter_data[quarter] = float(value_str)
                except ValueError:
                    pass

        if quarter_data and len(quarter_data) >= 2:
            chart_suggestions.append(
                {
                    "type": "bar",
                    "title": "분기별 실적",  # Quarterly Performance
                    "data": quarter_data,
                    "x_label": "분기",  # Quarter
                    "y_label": "매출 (원)",  # Sales (KRW)
                }
            )

        # If we couldn't extract real data, provide default charts
        if not chart_suggestions:
            print("Using default chart suggestions")
            chart_suggestions = [
                {
                    "type": "bar",
                    "title": "분기별 카테고리 매출",  # Quarterly Sales by Category
                    "data": {
                        "전자제품": 425000,  # Electronics
                        "의류": 310000,  # Clothing
                        "가정용품": 275000,  # Home Goods
                        "식품 및 음료": 190000,  # Food & Beverage
                        "뷰티": 150000,  # Beauty
                    },
                    "x_label": "제품 카테고리",  # Product Category
                    "y_label": "매출 (원)",  # Sales (KRW)
                },
                {
                    "type": "line",
                    "title": "월별 매출 추이",  # Monthly Sales Trend
                    "data": {
                        "1월": 120000,  # January
                        "2월": 115000,  # February
                        "3월": 130000,  # March
                        "4월": 140000,  # April
                        "5월": 135000,  # May
                        "6월": 155000,  # June
                    },
                    "x_label": "월",  # Month
                    "y_label": "매출 (원)",  # Sales (KRW)
                },
                {
                    "type": "pie",
                    "title": "지역별 매출 분포",  # Sales Distribution by Region
                    "data": {
                        "북부": 30,  # North
                        "남부": 25,  # South
                        "동부": 20,  # East
                        "서부": 15,  # West
                        "중부": 10,  # Central
                    },
                },
            ]

        return chart_suggestions
