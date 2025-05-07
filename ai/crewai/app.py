import os
from pathlib import Path

from dotenv import load_dotenv
from langchain_openai import ChatOpenAI

from crewai import Agent, Crew, Process, Task

# Load environment variables
load_dotenv()

# Set up the LLM
llm = ChatOpenAI(model="gpt-4o")

# Define the agents with specific roles
data_analyst = Agent(
    role="Data Analyst",
    goal="Analyze data and extract key insights for business reports",
    backstory="""You are an expert data analyst with years of experience in business intelligence.
    Your specialty is extracting meaningful insights from complex data and presenting them in a clear,
    actionable format. You have a talent for identifying trends and patterns that others might miss.""",
    verbose=True,
    llm=llm,
    allow_delegation=True,
)

content_writer = Agent(
    role="Content Writer",
    goal="Create compelling, clear, and concise content for business reports",
    backstory="""You are a skilled business writer with expertise in creating executive summaries,
    detailed analysis sections, and actionable recommendations. Your writing is known for being
    clear, professional, and tailored to the target audience.""",
    verbose=True,
    llm=llm,
    allow_delegation=True,
)

visual_designer = Agent(
    role="Visual Designer",
    goal="Create visually appealing and informative charts, graphs, and slide layouts",
    backstory="""You are a talented visual designer specializing in data visualization and
    presentation design. You know how to transform complex data into clear, compelling visuals
    that enhance understanding and engagement. You have a strong sense of design principles
    and know how to create professional business presentations.""",
    verbose=True,
    llm=llm,
    allow_delegation=True,
)

quality_reviewer = Agent(
    role="Quality Reviewer",
    goal="Ensure accuracy, consistency, and professionalism of the final report",
    backstory="""You are a meticulous reviewer with an eye for detail. You ensure that all
    business reports meet the highest standards of quality, accuracy, and professionalism.
    You check for consistency in messaging, formatting, and data presentation, and provide
    constructive feedback for improvements.""",
    verbose=True,
    llm=llm,
    allow_delegation=True,
)


# Define the tasks for each agent
def create_report_tasks(report_topic, data_sources=None):
    """
    Create tasks for generating a business report on the specified topic

    Args:
        report_topic (str): The main topic of the business report
        data_sources (list, optional): List of data sources to analyze

    Returns:
        list: List of Task objects
    """
    # Default data sources if none provided
    if data_sources is None:
        data_sources = [
            "Sample quarterly sales data",
            "Market trend analysis",
            "Customer feedback summary",
        ]

    # Data analysis task
    data_analysis_task = Task(
        description=f"""Analyze the following data sources for the {report_topic} report: {', '.join(data_sources)}.
        Identify key trends, patterns, and insights that would be valuable for a business report.
        Prepare a summary of your findings with specific data points that should be highlighted in the report.
        Suggest 3-5 key metrics or KPIs that should be visualized in charts or graphs.""",
        expected_output="A detailed analysis document with key insights, trends, and recommendations for data visualization",
        agent=data_analyst,
    )

    # Content creation task
    content_creation_task = Task(
        description=f"""Based on the data analysis provided, create the textual content for a business report on {report_topic}.
        The report should include:
        1. An executive summary (1-2 paragraphs)
        2. Introduction to the topic and its importance
        3. Key findings section with detailed analysis
        4. Strategic recommendations based on the data
        5. Conclusion summarizing the main points
        
        Ensure the content is professional, concise, and actionable. Use business-appropriate language and structure.""",
        expected_output="Complete text content for each section of the business report",
        agent=content_writer,
        context=[data_analysis_task],
    )

    # Visual design task
    visual_design_task = Task(
        description=f"""Create visual elements for the {report_topic} business report based on the data analysis and content provided.
        Design the following:
        1. Overall slide template and color scheme appropriate for a professional business presentation
        2. Data visualizations (charts, graphs) for the key metrics identified in the data analysis
        3. Layout recommendations for each slide, including placement of text, visuals, and any other elements
        4. Any additional visual elements that would enhance the presentation
        
        Provide detailed descriptions of each visual element and how it should be implemented in the final PowerPoint.""",
        expected_output="Detailed visual design specifications including chart types, layouts, and design elements",
        agent=visual_designer,
        context=[data_analysis_task, content_creation_task],
    )

    # Quality review task
    quality_review_task = Task(
        description=f"""Review all elements of the {report_topic} business report for quality, consistency, and professionalism.
        Check for:
        1. Accuracy of data and insights
        2. Clarity and professionalism of written content
        3. Effectiveness and appropriateness of visual elements
        4. Overall flow and structure of the presentation
        5. Alignment with business reporting best practices
        
        Provide specific feedback for improvements and ensure the final report meets high professional standards.""",
        expected_output="Comprehensive review with specific feedback and final approval of the business report",
        agent=quality_reviewer,
        context=[data_analysis_task, content_creation_task, visual_design_task],
    )

    return [
        data_analysis_task,
        content_creation_task,
        visual_design_task,
        quality_review_task,
    ]


# Function to generate the PowerPoint presentation
def generate_presentation(report_topic, crew_results):
    """
    Generate a PowerPoint presentation based on the crew's results

    Args:
        report_topic (str): The main topic of the business report
        crew_results (dict): Results from the crew's work

    Returns:
        str: Path to the generated PowerPoint file
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        # Create a new presentation
        prs = Presentation()

        # Extract content from crew results
        data_analysis = crew_results[0]
        content = crew_results[1]
        visual_design = crew_results[2]
        quality_review = crew_results[3]

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = f"{report_topic} Business Report"
        subtitle.text = "Generated by CrewAI Multi-Agent System"

        # Executive Summary slide
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]

        title.text = "Executive Summary"
        content_placeholder.text = (
            content.split("Executive Summary")[1].split("Introduction")[0].strip()
        )

        # Introduction slide
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]

        title.text = "Introduction"
        intro_text = content.split("Introduction")[1].split("Key Findings")[0].strip()
        content_placeholder.text = intro_text

        # Key Findings slides (multiple slides)
        findings_text = (
            content.split("Key Findings")[1]
            .split("Strategic Recommendations")[0]
            .strip()
        )
        findings_paragraphs = findings_text.split("\n\n")

        for i, paragraph in enumerate(findings_paragraphs):
            if paragraph.strip():
                slide = prs.slides.add_slide(content_slide_layout)
                title = slide.shapes.title
                content_placeholder = slide.placeholders[1]

                title.text = f"Key Findings {i+1}" if i > 0 else "Key Findings"
                content_placeholder.text = paragraph

        # Recommendations slide
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]

        title.text = "Strategic Recommendations"
        recommendations_text = (
            content.split("Strategic Recommendations")[1].split("Conclusion")[0].strip()
        )
        content_placeholder.text = recommendations_text

        # Conclusion slide
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]

        title.text = "Conclusion"
        conclusion_text = content.split("Conclusion")[1].strip()
        content_placeholder.text = conclusion_text

        # Save the presentation
        output_file = f"{report_topic.replace(' ', '_')}_Business_Report.pptx"
        prs.save(output_file)
        return output_file

    except Exception as e:
        print(f"Error generating presentation: {e}")
        return None


# Main function to run the report generation process
def generate_business_report(report_topic, data_sources=None):
    """
    Generate a complete business report on the specified topic

    Args:
        report_topic (str): The main topic of the business report
        data_sources (list, optional): List of data sources to analyze

    Returns:
        str: Path to the generated PowerPoint file
    """
    # Create tasks
    tasks = create_report_tasks(report_topic, data_sources)

    # Create the crew with all agents
    crew = Crew(
        agents=[data_analyst, content_writer, visual_designer, quality_reviewer],
        tasks=tasks,
        verbose=True,
        process=Process.sequential,
    )

    # Start the crew's work
    print(f"\nGenerating business report on: {report_topic}\n")
    results = crew.kickoff()

    # Generate the PowerPoint presentation
    print("\nCreating PowerPoint presentation...\n")
    presentation_path = generate_presentation(report_topic, results)

    if presentation_path:
        print(f"\nBusiness report successfully generated: {presentation_path}\n")
    else:
        print("\nFailed to generate PowerPoint presentation.\n")

    return presentation_path


# Example usage
if __name__ == "__main__":
    # Example topic and data sources
    topic = "Quarterly Sales Performance"
    sources = [
        "Q1-Q2 2023 Sales Data",
        "Competitor Analysis Report",
        "Customer Satisfaction Survey Results",
        "Market Trend Analysis",
    ]

    # Generate the report
    generate_business_report(topic, sources)
